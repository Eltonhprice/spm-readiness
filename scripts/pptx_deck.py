# scripts/pptx_deck.py
"""Generate a PowerPoint SPM Readiness leadership deck from metrics.json."""
import argparse
import io
import json
import math
import os
import tempfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from scripts.scoring import overall_score as _overall_score

# ── Palette ────────────────────────────────────────────────────────────────────
PURPLE    = RGBColor(0xA1, 0x00, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
GREY_LT   = RGBColor(0xF9, 0xF5, 0xFF)
GREY_MID  = RGBColor(0xE9, 0xD5, 0xFF)
GREY_TEXT = RGBColor(0x66, 0x66, 0x66)

RAG_RGB = {
    "green":         RGBColor(0x22, 0xC5, 0x5E),
    "amber":         RGBColor(0xF5, 0x9E, 0x0B),
    "red":           RGBColor(0xEF, 0x44, 0x44),
    "not_collected": RGBColor(0x9C, 0xA3, 0xAF),
}
RAG_HEX = {
    "green":  "#22c55e",
    "amber":  "#f59e0b",
    "red":    "#ef4444",
    "not_collected": "#9ca3af",
}

_MODULE_LABELS = {
    "demand":     "Demand Management",
    "ppm":        "Project Portfolio",
    "resource":   "Resource Mgmt",
    "financial":  "Financial Mgmt",
    "agile":      "Agile Development",
    "apm":        "APM",
    "innovation": "Innovation",
    "csdm":       "CSDM/CMDB Health",
}
_DIMS = ["activation", "data_volume", "data_completeness", "process_adoption", "integration"]
_DIM_SHORT = {
    "activation":        "Activation",
    "data_volume":       "Data Vol.",
    "data_completeness": "Completeness",
    "process_adoption":  "Adoption",
    "integration":       "Integration",
}

# ── Slide size (widescreen 13.33" × 7.5") ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _txbox(slide, x, y, w, h, text, size=12, bold=False, color=None,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or DARK
    return tb


def _rag_label(score):
    if score is None:
        return "not_collected"
    return "green" if score >= 70 else "amber" if score >= 40 else "red"


def _purple_bar(slide, y_top=Inches(0), height=Inches(0.08)):
    _rect(slide, 0, y_top, W, height, fill=PURPLE)


def _header_band(slide, title, subtitle=""):
    _rect(slide, 0, 0, W, Inches(1.2), fill=PURPLE)
    _txbox(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5),
           title, size=24, bold=True, color=WHITE)
    if subtitle:
        _txbox(slide, Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.4),
               subtitle, size=12, color=RGBColor(0xE9, 0xD5, 0xFF))


def _footer(slide, date="", mode="rde"):
    y = H - Inches(0.35)
    _rect(slide, 0, y, W, Inches(0.35), fill=PURPLE)
    label = f"Accenture  ·{mode.upper()}  ·  SPM Readiness Assessment  ·  {date}"
    _txbox(slide, Inches(0.4), y + Inches(0.05), Inches(12), Inches(0.25),
           label, size=9, color=RGBColor(0xE9, 0xD5, 0xFF))


# ── Radar chart (matplotlib → PNG → embed) ────────────────────────────────────
def _radar_png(labels, values, colors, size_px=480):
    n  = len(labels)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
    angles += angles[:1]

    vals = [v / 100.0 for v in values]
    vals += vals[:1]

    fig, ax = plt.subplots(figsize=(size_px/96, size_px/96), dpi=96,
                           subplot_kw=dict(polar=True))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")

    # Grid rings
    ax.set_ylim(0, 1)
    ax.set_yticks([0.25, 0.5, 0.75, 1.0])
    ax.set_yticklabels(["25%", "50%", "75%", "100%"], size=7, color="#aaa")
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, size=8, color="#333")
    ax.grid(color="#e5e7eb", linewidth=0.8)
    ax.spines["polar"].set_color("#e5e7eb")

    # Fill
    ax.fill(angles, vals, color="#A100FF", alpha=0.15)
    ax.plot(angles, vals, color="#A100FF", linewidth=2)

    # Dots per module colored by RAG
    for i, (ang, val, col) in enumerate(zip(angles[:-1], vals[:-1], colors)):
        ax.plot(ang, val, "o", color=col, markersize=6, zorder=5)

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=96, bbox_inches="tight",
                transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Slides ────────────────────────────────────────────────────────────────────
def _slide_cover(prs, client, date, overall, mode):
    sl = _blank(prs)
    # Purple left accent bar
    _rect(sl, 0, 0, Inches(0.18), H, fill=PURPLE)
    # Top accent line
    _rect(sl, Inches(0.18), 0, W - Inches(0.18), Inches(0.08), fill=PURPLE)

    # Badge
    _txbox(sl, Inches(0.5), Inches(0.35), Inches(12), Inches(0.4),
           f"ACCENTURE  ·  {mode.upper()}  ·  SPM READINESS ASSESSMENT  ·  AS-IS",
           size=9, bold=True, color=PURPLE)

    # Client name
    _txbox(sl, Inches(0.5), Inches(0.8), Inches(10), Inches(1.0),
           client, size=36, bold=True, color=DARK)

    # Date
    _txbox(sl, Inches(0.5), Inches(1.8), Inches(6), Inches(0.4),
           date, size=13, color=GREY_TEXT)

    # Score label
    _txbox(sl, Inches(0.5), Inches(2.6), Inches(6), Inches(0.4),
           "OVERALL SPM READINESS SCORE", size=10, bold=True,
           color=GREY_TEXT)

    # Big score
    rag = _rag_label(overall)
    score_str = f"{overall}%" if overall is not None else "—"
    _txbox(sl, Inches(0.45), Inches(2.9), Inches(6), Inches(2.2),
           score_str, size=96, bold=True,
           color=RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF)))

    # RAG legend
    legend = "Green ≥ 70%  ·  Amber 40–69%  ·  Red < 40%"
    _txbox(sl, Inches(0.5), Inches(5.4), Inches(8), Inches(0.4),
           legend, size=10, color=GREY_TEXT)

    # Footer bar
    _footer(sl, date, mode)
    return sl


def _slide_bar(prs, scores, overall, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Readiness at a Glance",
                 "Module scores  ·  Green ≥70%  ·  Amber 40–69%  ·  Red <40%  ·  Grey = not collected")

    # Layout constants
    label_x   = Inches(0.3)
    label_w   = Inches(2.3)
    bar_x     = Inches(2.75)
    bar_max_w = Inches(9.2)
    score_x   = Inches(12.1)
    score_w   = Inches(1.0)
    row_h     = Inches(0.52)
    gap       = Inches(0.06)
    y0        = Inches(1.42)

    for i, (key, label) in enumerate(_MODULE_LABELS.items()):
        ms  = scores.get(key, {}).get("module_score")
        rag = _rag_label(ms)
        col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))
        y   = y0 + i * (row_h + gap)

        # Module label (right-aligned feel via wide box)
        _txbox(sl, label_x, y + Inches(0.1), label_w, Inches(0.32),
               label, size=10, color=DARK)

        # Bar track (grey background)
        _rect(sl, bar_x, y + Inches(0.1), bar_max_w, Inches(0.32),
              fill=RGBColor(0xE9, 0xE9, 0xEB))

        # Bar fill — proportional to score
        fill_w = bar_max_w * (ms / 100.0) if ms is not None else 0
        if fill_w > 0:
            _rect(sl, bar_x, y + Inches(0.1), fill_w, Inches(0.32), fill=col)

        # Score label
        score_txt = f"{ms}%" if ms is not None else "—"
        _txbox(sl, score_x, y + Inches(0.08), score_w, Inches(0.36),
               score_txt, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # Overall score tile at bottom
    rag_o = _rag_label(overall)
    col_o = RAG_RGB.get(rag_o, RGBColor(0x9C, 0xA3, 0xAF))
    y_ov  = y0 + len(_MODULE_LABELS) * (row_h + gap) + Inches(0.12)
    _rect(sl, bar_x, y_ov, bar_max_w + score_w + Inches(0.15), Inches(0.55), fill=PURPLE)
    _txbox(sl, bar_x + Inches(0.2), y_ov + Inches(0.1), Inches(6), Inches(0.35),
           "OVERALL SPM READINESS", size=11, bold=True, color=WHITE)
    _txbox(sl, bar_x + bar_max_w - Inches(1.8), y_ov + Inches(0.1), Inches(2.0), Inches(0.35),
           f"{overall}%" if overall is not None else "—",
           size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    _footer(sl, date, mode)
    return sl


def _slide_scorecard(prs, scores, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Module Readiness Scorecard",
                 "8 modules × 5 dimensions  ·  RAG thresholds: Green ≥70%  ·  Amber 40–69%  ·  Red <40%")

    col_w  = [Inches(2.4)] + [Inches(1.4)] * 5 + [Inches(0.85), Inches(0.9)]
    row_h  = Inches(0.48)
    x0, y0 = Inches(0.25), Inches(1.35)
    headers = ["Module"] + [_DIM_SHORT[d] for d in _DIMS] + ["Score", "RAG"]

    total_w = sum(col_w)
    rows    = 1 + len(_MODULE_LABELS)
    tbl = sl.shapes.add_table(rows, len(headers), x0, y0, total_w,
                              row_h * rows).table

    # Header row
    for ci, (hdr, cw) in enumerate(zip(headers, col_w)):
        tbl.columns[ci].width = cw
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE

    for ri, (mod_key, mod_label) in enumerate(_MODULE_LABELS.items(), start=1):
        s      = scores.get(mod_key, {})
        ms     = s.get("module_score")
        mr     = s.get("rag", "not_collected")
        row_bg = GREY_LT if ri % 2 == 0 else WHITE

        tbl.rows[ri].height = row_h
        vals_row = [mod_label] + [s.get(d) for d in _DIMS] + [ms, mr.upper()]

        for ci, val in enumerate(vals_row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()

            if ci == 0:
                cell.fill.fore_color.rgb = row_bg
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = DARK
            elif ci <= 5:
                score_val = val
                rag = _rag_label(score_val) if score_val is not None else "not_collected"
                cell.fill.fore_color.rgb = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))
                display = f"{score_val}%" if score_val is not None else "—"
                cell.text = display
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE
            elif ci == 6:
                cell.fill.fore_color.rgb = row_bg
                cell.text = f"{ms}%" if ms is not None else "—"
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                col_ms = RAG_RGB.get(_rag_label(ms), RGBColor(0x9C, 0xA3, 0xAF))
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = col_ms
            else:
                col_rag = RAG_RGB.get(mr, RGBColor(0x9C, 0xA3, 0xAF))
                cell.fill.fore_color.rgb = col_rag
                cell.text = mr.upper()
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(8); run.font.bold = True; run.font.color.rgb = WHITE

    _footer(sl, date, mode)
    return sl


def _fmt(val, is_pct=False, invert=False):
    if val is None:
        return "—"
    if invert:
        val = 100.0 - float(val)
    if is_pct:
        return f"{round(float(val))}%"
    return str(val)


# Key metric inputs per module: (row_label, metric_key, is_pct, invert, dimension_label)
_MODULE_INPUTS = {
    "demand": [
        ("Total demands",           "total",                    False, False, "Data Volume"),
        ("Linked to project",       "linked_to_project_pct",    True,  False, "Integration"),
        ("Portfolio link",          "demand_with_portfolio_pct",True,  False, "Integration"),
        ("Program link",            "demand_with_program_pct",  True,  False, "Integration"),
        ("With approval",           "with_approval_pct",        True,  False, "Process Adoption"),
        ("Lifecycle throughput",    "demand_throughput_pct",    True,  False, "Process Adoption"),
        ("Workbench activity (14d)","demand_reviewed_14d_pct",  True,  False, "Process Adoption"),
        ("Stale >60d",              "stale_60d_pct",            True,  True,  "Process Adoption"),
        ("Priority set",            "demand_priority_set_pct",  True,  False, "Completeness"),
        ("Has owner",               "no_owner_pct",             True,  True,  "Completeness"),
    ],
    "ppm": [
        ("Total projects",          "total",                    False, False, "Data Volume"),
        ("Shell projects",          "shell_project_pct",        True,  False, "Completeness"),
        ("Has owner",               "no_owner_pct",             True,  True,  "Completeness"),
        ("Grouped in program",      "with_program_pct",         True,  False, "Integration"),
        ("With approval",           "with_approval_pct",        True,  False, "Process Adoption"),
        ("Stale >30d",              "stale_30d_pct",            True,  True,  "Process Adoption"),
    ],
    "resource": [
        ("Total resource plans","total",                           False, False, "Data Volume"),
        ("Named resource",      "resource_plan_named_pct",         True,  False, "Completeness"),
        ("Linked to project",   "linked_to_project_pct",           True,  False, "Integration"),
        ("Timesheet coverage",  "timesheet_coverage_pct",          True,  False, "Process Adoption"),
        ("Req. stale >30d",     "resource_requests_stale_30d_pct", True,  True,  "Process Adoption"),
    ],
    "financial": [
        ("With financials",     "projects_with_financials_pct", True, False, "Integration"),
        ("With cost plan",      "projects_with_cost_plan_pct",  True, False, "Completeness"),
        ("With budget plan",    "projects_with_budget_plan_pct",True, False, "Completeness"),
    ],
    "agile": [
        ("Total stories",           "total_stories",            False, False, "Data Volume"),
        ("Stories with sprint",     "no_sprint_pct",            True,  True,  "Completeness"),
        ("Stories with team",       "no_team_pct",              True,  True,  "Completeness"),
        ("Completed sprints",       "completed_sprint_count",   False, False, "Process Adoption"),
        ("Backlog stale >45d",      "backlog_stale_45d_pct",    True,  True,  "Process Adoption"),
    ],
    "apm": [
        ("Total applications",  "total",                    False, False, "Data Volume"),
        ("With lifecycle stage","with_lifecycle_stage_pct", True,  False, "Completeness"),
        ("Has owner",           "with_owner_pct",           True,  False, "Completeness"),
        ("Linked to CMDB",      "with_cmdb_link_pct",       True,  False, "Integration"),
    ],
    "innovation": [
        ("Total ideas",         "total",                    False, False, "Data Volume"),
        ("Has owner",           "no_owner_pct",             True,  True,  "Completeness"),
        ("Linked to project",   "linked_to_demand_or_project_pct", True, False, "Integration"),
    ],
    "csdm": [
        ("Total CIs",           "total_ci",                 False, False, "Data Volume"),
        ("Total services",      "total_services",           False, False, "Data Volume"),
        ("CI: operational status set","ci_with_operational_status_pct", True, False, "Completeness"),
        ("CI: has owner",       "ci_with_owner_pct",        True,  False, "Completeness"),
        ("CI: has support group","ci_with_support_group_pct",True,  False, "Completeness"),
        ("CI: auto-discovered", "ci_discovered_pct",        True,  False, "Process Adoption"),
        ("Services with owner", "services_with_owner_pct",  True,  False, "Process Adoption"),
        ("Total relationships", "total_relationships",      False, False, "Integration"),
    ],
}


def _slide_scoring_basis(prs, metrics, scores, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Scoring Basis",
                 "Key metric inputs that drove each module's dimension scores")

    mods = metrics.get("modules", {})
    keys = list(_MODULE_LABELS.keys())
    left_keys  = keys[:4]   # demand, ppm, resource, financial
    right_keys = keys[4:]   # agile, apm, innovation, csdm

    col_x     = [Inches(0.25), Inches(6.8)]
    col_w     = Inches(6.4)
    y_start   = Inches(1.38)
    row_h     = Inches(0.3)
    hdr_h     = Inches(0.36)
    gap_mod   = Inches(0.1)

    for col_i, col_keys in enumerate([left_keys, right_keys]):
        x  = col_x[col_i]
        y  = y_start

        for mod_key in col_keys:
            label  = _MODULE_LABELS[mod_key]
            ms     = scores.get(mod_key, {}).get("module_score")
            rag    = _rag_label(ms)
            col    = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))
            mod    = mods.get(mod_key, {})
            inputs = _MODULE_INPUTS.get(mod_key, [])

            # Module header bar
            _rect(sl, x, y, col_w, hdr_h, fill=PURPLE)
            _txbox(sl, x + Inches(0.12), y + Inches(0.05), Inches(4.0), Inches(0.26),
                   label, size=10, bold=True, color=WHITE)
            score_str = f"{ms}%" if ms is not None else "Not collected"
            _txbox(sl, x + Inches(4.1), y + Inches(0.05), Inches(2.1), Inches(0.26),
                   score_str, size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
            y += hdr_h

            if ms is None and not mod.get("plugin_active", False):
                # Not collected — single grey row
                _rect(sl, x, y, col_w, row_h, fill=GREY_LT)
                _txbox(sl, x + Inches(0.15), y + Inches(0.04), col_w - Inches(0.2), Inches(0.22),
                       "Plugin not active — not collected", size=8,
                       color=GREY_TEXT, italic=True)
                y += row_h
            else:
                for row_i, (row_lbl, metric_key, is_pct, invert, dim_lbl) in enumerate(inputs):
                    raw = mod.get(metric_key)
                    val_str = _fmt(raw, is_pct=is_pct, invert=invert)
                    bg = GREY_LT if row_i % 2 == 0 else WHITE
                    _rect(sl, x, y, col_w, row_h, fill=bg)
                    _txbox(sl, x + Inches(0.15), y + Inches(0.04),
                           Inches(2.8), Inches(0.22),
                           row_lbl, size=8, color=DARK)
                    _txbox(sl, x + Inches(3.0), y + Inches(0.04),
                           Inches(1.5), Inches(0.22),
                           val_str, size=8, bold=(raw is not None), color=DARK,
                           align=PP_ALIGN.CENTER)
                    _txbox(sl, x + Inches(4.6), y + Inches(0.04),
                           Inches(1.7), Inches(0.22),
                           dim_lbl, size=7, color=GREY_TEXT, italic=True,
                           align=PP_ALIGN.RIGHT)
                    y += row_h

            y += gap_mod

    _footer(sl, date, mode)
    return sl


def _slide_governance(prs, metrics, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Governance & Process Adoption",
                 "Approval, timesheet, status reporting, and scoring model signals")

    gov   = metrics.get("governance", {})
    ts    = gov.get("timesheets", {})
    app   = gov.get("approvals", {})
    sr    = gov.get("status_reports", {})
    sm    = gov.get("scoring_models", {})
    pa    = metrics.get("pa_adoption", {})
    roles = metrics.get("roles", {})

    signals = [
        ("Active Timesheet Periods",                    ts.get("active_periods"),                 ts.get("collected")),
        ("Timesheet Entries (total)",                   ts.get("total_entries"),                  ts.get("collected")),
        ("Demands with Approvals",                      app.get("demand_records_with_approvals"),  app.get("collected")),
        ("Projects with Approvals",                     app.get("project_records_with_approvals"), app.get("collected")),
        ("Project Status Reports",                      sr.get("total"),                           sr.get("collected")),
        ("PA Scorecards (Resource Mgmt Dashboard)*",    pa.get("scorecard_count"),                 pa.get("collected")),
        ("Portfolio Scoring Criteria",                  sm.get("criteria_count"),                  sm.get("collected")),
        ("Records with Portfolio Score",                sm.get("scored_records"),                  sm.get("collected")),
    ]

    col_w = [Inches(4.5), Inches(2.2), Inches(2.5)]
    row_h = Inches(0.5)
    x0, y0 = Inches(0.5), Inches(1.4)
    tbl = sl.shapes.add_table(len(signals) + 1, 3, x0, y0,
                              sum(col_w), row_h * (len(signals) + 1)).table
    for ci, (hdr, cw) in enumerate(zip(["Signal", "Value", "Status"], col_w)):
        tbl.columns[ci].width = cw
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = WHITE

    for ri, (label, val, collected) in enumerate(signals, start=1):
        bg = GREY_LT if ri % 2 == 0 else WHITE
        tbl.rows[ri].height = row_h
        for ci in range(3):
            cell = tbl.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tbl.cell(ri, 0).text = label
        tbl.cell(ri, 1).text = str(val) if val is not None else "—"
        tbl.cell(ri, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        status_txt = "Collected" if collected else "Not collected"
        status_col = RAG_RGB["green"] if collected else RGBColor(0x9C, 0xA3, 0xAF)
        status_cell = tbl.cell(ri, 2)
        status_cell.fill.solid(); status_cell.fill.fore_color.rgb = status_col
        status_cell.text = status_txt
        p = status_cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE
        for ci in [0, 1]:
            p = tbl.cell(ri, ci).text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(10); run.font.color.rgb = DARK

    # Role counts side panel
    x1 = Inches(7.7)
    _rect(sl, x1, Inches(1.4), Inches(5.3), Inches(5.4), fill=GREY_LT, line=GREY_MID)
    _txbox(sl, x1 + Inches(0.2), Inches(1.5), Inches(5), Inches(0.4),
           "SPM ROLE ASSIGNMENTS", size=9, bold=True, color=PURPLE)
    role_labels = {
        "portfolio_manager": "Portfolio Manager",
        "project_manager":   "Project Manager",
        "resource_manager":  "Resource Manager",
        "financial_analyst": "Financial Analyst",
        "it_demand_manager": "IT Demand Manager",
    }
    for i, (key, label) in enumerate(role_labels.items()):
        y = Inches(2.05) + i * Inches(0.7)
        count = roles.get(key)
        _txbox(sl, x1 + Inches(0.2), y, Inches(3.5), Inches(0.35),
               label, size=10, color=DARK)
        _txbox(sl, x1 + Inches(3.8), y, Inches(1.3), Inches(0.35),
               str(count) if count is not None else "—", size=13, bold=True,
               color=PURPLE, align=PP_ALIGN.RIGHT)

    _txbox(sl, Inches(0.5), Inches(6.85), Inches(7.0), Inches(0.28),
           "* PA Scorecards: Performance Analytics plugin required. Scorecard count is an aggregate signal — "
           "it cannot confirm whether every individual resource manager is performing regular reviews.",
           size=8, color=GREY_MID)

    _footer(sl, date, mode)
    return sl


def _slide_staleness(prs, metrics, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Staleness & Data Quality",
                 "Records not updated within module threshold — incorporated into process-adoption scores")

    dq    = metrics.get("data_quality", {})
    agile = metrics.get("modules", {}).get("agile", {})
    res   = metrics.get("modules", {}).get("resource", {})

    rows = [
        ("Projects not updated (active)",        dq.get("projects_stale_30d"),     dq.get("projects_stale_30d_pct"),     "30 days", "PPM → Process Adoption"),
        ("Demands not updated (open)",            dq.get("demands_stale_60d"),      dq.get("demands_stale_60d_pct"),      "60 days", "Demand → Process Adoption"),
        ("Backlog stories not updated",           None,                             agile.get("backlog_stale_45d_pct"),   "45 days", "Agile → Process Adoption"),
        ("Open resource requests not updated",    None,                             res.get("resource_requests_stale_30d_pct"), "30 days", "Resource → Process Adoption"),
    ]

    col_w  = [Inches(3.6), Inches(1.2), Inches(1.4), Inches(1.4), Inches(3.4)]
    hdrs   = ["Indicator", "Count", "% Stale", "Threshold", "Scores Into"]
    row_h  = Inches(0.55)
    x0, y0 = Inches(0.4), Inches(1.42)
    total_w = sum(col_w)

    tbl = sl.shapes.add_table(len(rows) + 1, 5, x0, y0,
                               total_w, row_h * (len(rows) + 1)).table
    for ci, (hdr, cw) in enumerate(zip(hdrs, col_w)):
        tbl.columns[ci].width = cw
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = WHITE

    for ri, (label, count, pct, threshold, scores_into) in enumerate(rows, start=1):
        bg = GREY_LT if ri % 2 == 0 else WHITE
        tbl.rows[ri].height = row_h
        for ci in range(5):
            tbl.cell(ri, ci).fill.solid()
            tbl.cell(ri, ci).fill.fore_color.rgb = bg

        count_str = str(count) if count is not None else "—"
        pct_str   = f"{pct:.1f}%" if pct is not None else "—"

        # RAG-colour the % cell: >50% stale = red, >20% = amber, else green
        if pct is not None:
            rag_col = RAG_RGB["red"] if pct > 50 else (RAG_RGB["amber"] if pct > 20 else RAG_RGB["green"])
        else:
            rag_col = RGBColor(0x9C, 0xA3, 0xAF)

        vals = [label, count_str, pct_str, threshold, scores_into]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
        for ci, (v, al) in enumerate(zip(vals, aligns)):
            cell = tbl.cell(ri, ci)
            cell.text = v
            p = cell.text_frame.paragraphs[0]
            p.alignment = al
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(10)
            run.font.color.rgb = DARK
            if ci == 2 and pct is not None:   # % column — bold + rag fill
                cell.fill.solid(); cell.fill.fore_color.rgb = rag_col
                run.font.color.rgb = WHITE; run.font.bold = True

    # Explanatory note below table
    note_y = y0 + row_h * (len(rows) + 1) + Inches(0.2)
    _txbox(sl, x0, note_y, total_w, Inches(0.5),
           "Note: staleness percentages are incorporated into each module's process-adoption dimension score — "
           "they are not reported as a standalone dimension. Thresholds reflect realistic update cadences per module type.",
           size=8, color=GREY_TEXT, italic=True)

    _footer(sl, date, mode)
    return sl


def _slide_findings(prs, findings, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Top Findings",
                 "Highest-priority signals from the readiness assessment")

    top5 = findings[:5]
    if not top5:
        _txbox(sl, Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
               "No significant findings generated from the available data.",
               size=12, color=GREY_TEXT)
    else:
        col_w = [Inches(0.9), Inches(2.3), Inches(7.3), Inches(2.4)]
        row_h = Inches(0.72)
        x0, y0 = Inches(0.25), Inches(1.35)
        tbl = sl.shapes.add_table(len(top5) + 1, 4, x0, y0,
                                  sum(col_w), row_h * (len(top5) + 1)).table
        for ci, (hdr, cw) in enumerate(zip(["ID", "Module", "Observation", "Significance"], col_w)):
            tbl.columns[ci].width = cw
            cell = tbl.cell(0, ci)
            cell.text = hdr
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE

        for ri, f in enumerate(top5, start=1):
            rag = f.get("rag", "not_collected")
            col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))
            bg  = GREY_LT if ri % 2 == 0 else WHITE
            tbl.rows[ri].height = row_h

            id_cell = tbl.cell(ri, 0)
            id_cell.fill.solid(); id_cell.fill.fore_color.rgb = col
            id_cell.text = f.get("id", "")
            p = id_cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE

            for ci, txt in [(1, f.get("module_label", f.get("module", ""))),
                            (2, f.get("observation", "")),
                            (3, f.get("significance", ""))]:
                cell = tbl.cell(ri, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.text = txt
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(9); run.font.color.rgb = DARK

    _footer(sl, date, mode)
    return sl


def _slide_next_steps(prs, mode, findings, date):
    sl = _blank(prs)
    if mode == "fde" and findings:
        _header_band(sl, "Priority Focus Areas",
                     "FDE · Areas for Discussion  ·  Derived from top Red/Amber findings")
        priority = [f for f in findings if f.get("rag") in ("red", "amber")][:5]
        if not priority:
            priority = findings[:5]

        col_w = [Inches(0.5), Inches(2.3), Inches(7.5), Inches(2.4)]
        row_h = Inches(0.72)
        x0, y0 = Inches(0.25), Inches(1.35)
        tbl = sl.shapes.add_table(len(priority) + 1, 4, x0, y0,
                                  sum(col_w), row_h * (len(priority) + 1)).table
        for ci, (hdr, cw) in enumerate(zip(["#", "Module", "Finding", "Priority"], col_w)):
            tbl.columns[ci].width = cw
            cell = tbl.cell(0, ci)
            cell.text = hdr
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE

        for ri, f in enumerate(priority, start=1):
            rag = f.get("rag", "not_collected")
            col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))
            bg  = GREY_LT if ri % 2 == 0 else WHITE
            tbl.rows[ri].height = row_h

            num_cell = tbl.cell(ri, 0)
            num_cell.fill.solid(); num_cell.fill.fore_color.rgb = PURPLE
            num_cell.text = str(ri)
            p = num_cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = WHITE

            for ci, txt in [(1, f.get("module_label", f.get("module", ""))),
                            (2, f.get("observation", ""))]:
                cell = tbl.cell(ri, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.text = txt
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(9); run.font.color.rgb = DARK

            pri_cell = tbl.cell(ri, 3)
            pri_cell.fill.solid(); pri_cell.fill.fore_color.rgb = col
            pri_cell.text = rag.upper()
            p = pri_cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE

        note = "Consultant to validate scope and sequencing with client."
        _txbox(sl, Inches(0.25), Inches(6.85), Inches(12), Inches(0.3),
               note, size=9, color=GREY_TEXT, italic=True)
    else:
        _header_band(sl, "Recommended Next Steps", "For Discussion — consultant to complete")
        _rect(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.4),
              fill=GREY_LT, line=GREY_MID)
        _txbox(sl, Inches(0.7), Inches(1.65), Inches(11.8), Inches(0.4),
               "FOR DISCUSSION", size=10, bold=True, color=PURPLE)
        _txbox(sl, Inches(0.7), Inches(2.1), Inches(11.8), Inches(0.8),
               "This slide is intentionally left as a structured placeholder.\n"
               "The consultant completes this section during the FDE conversation with the client.",
               size=11, color=GREY_TEXT)

        col_w = [Inches(0.6), Inches(5.5), Inches(2.5), Inches(3.5)]
        x0, y0 = Inches(0.5), Inches(4.1)
        tbl = sl.shapes.add_table(4, 4, x0, y0, sum(col_w), Inches(2.0)).table
        for ci, (hdr, cw) in enumerate(zip(["#", "Initiative", "Module", "Priority"], col_w)):
            tbl.columns[ci].width = cw
            cell = tbl.cell(0, ci)
            cell.text = hdr
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE
        for ri in range(1, 4):
            bg = GREY_LT if ri % 2 == 0 else WHITE
            tbl.cell(ri, 0).text = str(ri)
            for ci in range(4):
                cell = tbl.cell(ri, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                if ci > 0:
                    cell.text = "[Consultant to complete]"
                    p = cell.text_frame.paragraphs[0]
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.size = Pt(9); run.font.color.rgb = GREY_TEXT
                    run.font.italic = True

    _footer(sl, date, mode)
    return sl


_PLUGIN_IDS = {
    "demand":     "com.snc.sdlc.ppm_core (shared with PPM)",
    "ppm":        "com.snc.sdlc.ppm_core",
    "resource":   "com.snc.rm",
    "financial":  "com.snc.financial_mgmt",
    "agile":      "com.snc.agile",
    "apm":        "com.snc.apm",
    "innovation": "com.snc.innovation_mgmt",
    "csdm":       "cmdb_ci table (data-driven)",
}

_PRIMARY_TABLES = {
    "demand":     "pm_demand",
    "ppm":        "pm_project",
    "resource":   "pm_resource_plan",
    "financial":  "pm_project_financials",
    "agile":      "rm_story",
    "apm":        "apm_appl_now",
    "innovation": "innovation_idea",
    "csdm":       "cmdb_ci",
}

_COMPLETENESS_FIELDS = {
    "demand":     "demand_priority_set, assigned_to",
    "ppm":        "percent_complete, assigned_to, phase (shell check)",
    "resource":   "resource.name filled, actual vs planned hours",
    "financial":  "cost_plan linked, budget_plan linked, actuals present",
    "agile":      "sprint assigned, team assigned",
    "apm":        "lifecycle_stage, owned_by",
    "innovation": "assigned_to",
    "csdm":       "operational_status, owned_by, support_group, environment",
}

_ADOPTION_SIGNALS = {
    "demand":     "approval workflow rate + scoring model usage",
    "ppm":        "status report filed ≤30 days + approval workflow rate",
    "resource":   "timesheet entries logged (>200=100%, >50=60%, >0=20%)",
    "financial":  "project approvals logged (>10=100%, >2=60%, >0=20%)",
    "agile":      "completed sprints (>5=100%, >1=60%, >0=20%)",
    "apm":        "not assessed (no standard governance signal)",
    "innovation": "% ideas linked to demand or project",
    "csdm":       "% CIs discovered (not manually entered) + % services with owner",
}

_INTEGRATION_SIGNAL = {
    "demand":     "% demands linked to a project",
    "ppm":        "% projects grouped under a program",
    "resource":   "% resource plans linked to a project",
    "financial":  "% PPM projects with any financial record",
    "agile":      "% stories assigned to a team",
    "apm":        "% applications linked to a CMDB CI",
    "innovation": "% ideas linked to demand or project",
    "csdm":       "relationship density (total CMDB rels ÷ total CIs)",
}


def _slide_scorecard_explainer(prs, metrics, scores, date, mode):
    sl = _blank(prs)
    _header_band(sl, "How to Read the Scorecard",
                 "RAG thresholds · dimension weights · which modules are active on this instance")

    # ── Top-left: RAG thresholds (compact) ───────────────────────────────────
    _txbox(sl, Inches(0.3), Inches(1.38), Inches(3.8), Inches(0.24),
           "RAG THRESHOLDS", size=8, bold=True, color=GREY_TEXT)
    thresholds = [
        ("Green",    "Score ≥ 70%",                          RAG_RGB["green"]),
        ("Amber",    "40% ≤ score < 70%",                    RAG_RGB["amber"]),
        ("Red",      "Score < 40%",                          RAG_RGB["red"]),
        ("— (grey)", "Plugin inactive or no data collected", RGBColor(0x9C, 0xA3, 0xAF)),
    ]
    for i, (label, meaning, col) in enumerate(thresholds):
        y = Inches(1.64) + i * Inches(0.50)
        _rect(sl, Inches(0.3), y, Inches(1.2), Inches(0.40), fill=col)
        _txbox(sl, Inches(0.33), y + Inches(0.10), Inches(1.14), Inches(0.22),
               label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _rect(sl, Inches(1.55), y, Inches(4.35), Inches(0.40), fill=GREY_LT)
        _txbox(sl, Inches(1.65), y + Inches(0.10), Inches(4.15), Inches(0.22),
               meaning, size=9, color=DARK)

    # ── Top-right: dimension weights ──────────────────────────────────────────
    _txbox(sl, Inches(6.3), Inches(1.38), Inches(6.9), Inches(0.24),
           "DIMENSION WEIGHTS  (scored dims only — inactive dims excluded from average)",
           size=8, bold=True, color=GREY_TEXT)
    dims_w = [
        ("Activation",        0.20, "Plugin installed & active on instance"),
        ("Data Volume",       0.20, "Enough records to assess the module"),
        ("Data Completeness", 0.25, "Key fields populated across records"),
        ("Process Adoption",  0.25, "Governance mechanisms actively used"),
        ("Integration",       0.10, "Records linked across SPM modules"),
    ]
    bar_max = Inches(2.6)
    for i, (dim, w, desc) in enumerate(dims_w):
        y = Inches(1.64) + i * Inches(0.50)
        _rect(sl, Inches(6.3), y, Inches(6.9), Inches(0.40), fill=GREY_LT)
        _txbox(sl, Inches(6.42), y + Inches(0.04), Inches(1.9), Inches(0.20),
               dim, size=9, bold=True, color=DARK)
        _txbox(sl, Inches(6.42), y + Inches(0.22), Inches(4.2), Inches(0.16),
               desc, size=7, color=GREY_TEXT, italic=True)
        _rect(sl, Inches(10.75), y + Inches(0.06), bar_max, Inches(0.26),
              fill=RGBColor(0xE2, 0xE8, 0xF0))
        _rect(sl, Inches(10.75), y + Inches(0.06), bar_max * (w / 0.25), Inches(0.26),
              fill=PURPLE)
        _txbox(sl, Inches(13.45), y + Inches(0.07), Inches(0.7), Inches(0.22),
               f"{int(w*100)}%", size=9, bold=True, color=PURPLE, align=PP_ALIGN.RIGHT)

    # ── Bottom: per-module plugin / activation status table ───────────────────
    _txbox(sl, Inches(0.3), Inches(4.18), Inches(12.7), Inches(0.24),
           "MODULE STATUS ON THIS INSTANCE", size=8, bold=True, color=GREY_TEXT)

    # Table header
    hdr_y = Inches(4.44)
    _rect(sl, Inches(0.3), hdr_y, Inches(12.73), Inches(0.30), fill=PURPLE)
    hdrs = [("Module", 2.2), ("Plugin / Source", 3.1), ("Primary Table", 1.9),
            ("Records", 1.1), ("Score", 0.8), ("Status", 3.5)]
    hx = Inches(0.3)
    for hdr, cw in hdrs:
        _txbox(sl, hx + Inches(0.06), hdr_y + Inches(0.05), Inches(cw - 0.1), Inches(0.20),
               hdr, size=8, bold=True, color=WHITE)
        hx += Inches(cw)

    mods_data = metrics.get("modules", {})
    row_h = Inches(0.33)
    for ri, mod_key in enumerate(_MODULE_LABELS):
        mod  = mods_data.get(mod_key, {})
        ms   = scores.get(mod_key, {}).get("module_score")
        rag  = _rag_label(ms)
        plug = mod.get("plugin_active")

        # record count
        if mod_key == "csdm":
            n = mod.get("total_ci")
            unit = "CIs"
        elif mod_key == "agile":
            n = mod.get("total_stories")
            unit = "stories"
        else:
            n = mod.get("total")
            unit = "records"
        rec_str = f"{n:,} {unit}" if n else "—"

        # score cell
        score_str = f"{ms}%" if ms is not None else "—"
        score_col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))

        # status reason
        if ms is not None:
            status_str = f"Scored  [{rag.upper()}]"
            status_col = score_col
        elif not plug:
            status_str = "Plugin not active on instance"
            status_col = RGBColor(0x9C, 0xA3, 0xAF)
        elif not n:
            status_str = "No records collected"
            status_col = RGBColor(0x9C, 0xA3, 0xAF)
        else:
            status_str = "Not collected"
            status_col = RGBColor(0x9C, 0xA3, 0xAF)

        y = hdr_y + Inches(0.30) + ri * row_h
        bg = GREY_LT if ri % 2 == 0 else WHITE
        _rect(sl, Inches(0.3), y, Inches(12.73), row_h, fill=bg)

        vals = [
            (_MODULE_LABELS[mod_key],         2.2,  DARK,      False),
            (_PLUGIN_IDS.get(mod_key, "—"),   3.1,  GREY_TEXT, False),
            (_PRIMARY_TABLES.get(mod_key,"—"),1.9,  GREY_TEXT, False),
            (rec_str,                          1.1,  DARK,      True),
            (score_str,                        0.8,  score_col, True),
            (status_str,                       3.5,  status_col,False),
        ]
        vx = Inches(0.3)
        for val, cw, fc, bold in vals:
            _txbox(sl, vx + Inches(0.06), y + Inches(0.07), Inches(cw - 0.1), Inches(0.20),
                   val, size=8, bold=bold, color=fc)
            vx += Inches(cw)

    _footer(sl, date, mode)
    return sl


def _chip_row(sl, x, y, chip_color, text, size=8, text_color=None):
    """Draw one RAG chip + label on a single line. Returns y of next line."""
    chip_w, chip_h, line_h = Inches(0.18), Inches(0.16), Inches(0.22)
    if chip_color is not None:
        _rect(sl, x, y + Inches(0.03), chip_w, chip_h, fill=chip_color)
        _txbox(sl, x + chip_w + Inches(0.05), y, Inches(5.0), line_h,
               text, size=size, color=text_color or DARK)
    else:
        _txbox(sl, x + Inches(0.04), y, Inches(5.0), line_h,
               text, size=size - 1, color=GREY_TEXT, italic=True)
    return y + line_h


def _slide_scoring_basis_explainer(prs, metrics, scores, date, mode):
    sl = _blank(prs)
    _header_band(sl, "How Dimension Scores Are Derived",
                 "What each dimension measures · which table is queried · how the score is calculated")

    # Layout: 3 columns
    # Col A — Dimension name + weight  (x=0.25, w=2.5")
    # Col B — What it measures         (x=2.80, w=4.5")
    # Col C — How it scores            (x=7.35, w=5.75")
    COL_A_X, COL_A_W = Inches(0.25), Inches(2.5)
    COL_B_X, COL_B_W = Inches(2.80), Inches(4.5)
    COL_C_X, COL_C_W = Inches(7.35), Inches(5.75)
    row_h = Inches(1.12)
    y0    = Inches(1.35)

    # Table header bar
    _rect(sl, Inches(0.25), y0, Inches(12.85), Inches(0.30), fill=PURPLE)
    for txt, cx, cw in [
        ("Dimension  (weight)", COL_A_X, COL_A_W),
        ("What it measures  ·  source table", COL_B_X, COL_B_W),
        ("How it scores  —  visual threshold guide", COL_C_X, COL_C_W),
    ]:
        _txbox(sl, cx + Inches(0.08), y0 + Inches(0.06), cw - Inches(0.1), Inches(0.20),
               txt, size=8, bold=True, color=WHITE)

    # Vertical column dividers (light lines)
    for dx in [COL_B_X, COL_C_X]:
        _rect(sl, dx - Inches(0.04), y0, Inches(0.04), row_h * 5 + Inches(0.30),
              fill=GREY_MID)

    vol100 = {"demand": 50, "ppm": 20, "resource": 10, "financial": 10,
              "agile": 50, "apm": 10, "innovation": 5, "csdm": 500}
    vol60  = {"demand": 10, "ppm": 5,  "resource": 5,  "financial": 3,
              "agile": 10, "apm": 3,   "innovation": 2, "csdm": 50}

    # vol example strings — show 3 most relevant modules
    vol_green_ex = "≥20 projects · ≥50 stories · ≥500 CIs · ≥50 demands"
    vol_amber_ex = "≥5 projects  · ≥10 stories  · ≥50 CIs  · ≥10 demands"

    rows = [
        # (dim, pct, what_text, chip_lines)
        # chip_lines = list of (color_or_None, label)
        (
            "Activation", "20%",
            "Is the module's ServiceNow plugin installed and\n"
            "active on this instance?\n\n"
            "Each module maps to a unique plugin ID.\n"
            "Demand is an exception — it shares a plugin ID\n"
            "with PPM and also requires ≥1 demand record.",
            [
                (RAG_RGB["green"],           "Plugin active on instance  →  100%"),
                (RGBColor(0x9C,0xA3,0xAF),   "Plugin not active  →  Not Collected (—)"),
                (None, "Not Collected dims are excluded from the"),
                (None, "module score — they do not count as zero."),
            ],
        ),
        (
            "Data Volume", "20%",
            "How many records exist in the module's primary\n"
            "table? If there are too few records, the module\n"
            "cannot be meaningfully assessed.\n\n"
            "Primary tables:\n"
            "  PPM → pm_project\n"
            "  Agile → rm_story\n"
            "  CSDM → cmdb_ci\n"
            "  Demand → pm_demand    (others similar)",
            [
                (RAG_RGB["green"], f"≥ full threshold  →  100%"),
                (None,             f"  {vol_green_ex}"),
                (RAG_RGB["amber"], f"≥ partial threshold  →  60%"),
                (None,             f"  {vol_amber_ex}"),
                (RAG_RGB["red"],   "≥ 1 record  →  20%"),
                (RGBColor(0x9C,0xA3,0xAF), "0 records  →  Not Collected (—)"),
            ],
        ),
        (
            "Data Completeness", "25%",
            "What % of records have key fields populated?\n"
            "Each field is scored as a fill rate (0–100%).\n"
            "The dimension score = average of all field rates.\n\n"
            "Fields checked per module:\n"
            "  PPM: owner, % complete, phase (shell check)\n"
            "  Agile: sprint assigned, team assigned\n"
            "  CSDM: op. status, owner, support group, env\n"
            "  APM: lifecycle stage, owned_by\n"
            "  Others: see Scoring Basis slide",
            [
                (None, "Score = average of field fill rates  (0–100%)"),
                (None, ""),
                (None, "Null fields are excluded — not counted as zero."),
                (None, ""),
                (None, "Inverted fields: 100 − fill rate"),
                (None, "  (used when lower value = better quality,"),
                (None, "   e.g. shell-project % or no-owner %)"),
            ],
        ),
        (
            "Process Adoption", "25%",
            "Are people actually using the module's governance\n"
            "features — or is the plugin on but idle?\n\n"
            "Signals used by module:\n"
            "  PPM: status reports filed ≤30d + approval rate\n"
            "  Agile: number of completed sprints\n"
            "  Resource: timesheet entries logged\n"
            "  Financial: project approvals logged\n"
            "  Demand: approval rate + scoring model usage\n"
            "  CSDM: % discovered CIs + % services with owner\n"
            "  APM: not assessed (no standard signal)",
            [
                (RAG_RGB["green"], "High activity  →  100%"),
                (None,             "  >5 sprints · >200 timesheets · >10 approvals"),
                (RAG_RGB["amber"], "Moderate activity  →  60%"),
                (None,             "  >1 sprint · >50 timesheets · >2 approvals"),
                (RAG_RGB["red"],   "Any activity  →  20%"),
                (None,             ""),
                (None,             "0 activity  =  0%  (NOT Not Collected)"),
                (None,             "Absence of process is a real finding."),
            ],
        ),
        (
            "Integration", "10%",
            "Are this module's records linked to adjacent\n"
            "SPM modules, showing platform-wide usage?\n\n"
            "Linkage checked per module:\n"
            "  PPM: projects grouped under a program\n"
            "  Demand: demands linked to a project\n"
            "  Resource: resource plans linked to a project\n"
            "  Financial: projects with any financial record\n"
            "  Agile: stories assigned to a team\n"
            "  APM: applications linked to a CMDB CI\n"
            "  CSDM: relationship density (rels ÷ CIs)",
            [
                (None,             "Most modules: % of records with cross-module link"),
                (None,             "  0% – 100%  (higher = better integrated)"),
                (None,             ""),
                (None,             "CSDM uses relationship density instead:"),
                (RAG_RGB["green"], "≥ 2.0 relationships per CI  →  100%"),
                (RAG_RGB["amber"], "≥ 0.5 relationships per CI  →  60%"),
                (RAG_RGB["red"],   "< 0.5 relationships per CI  →  20%"),
            ],
        ),
    ]

    for ri, (dim, pct, what, chip_lines) in enumerate(rows):
        y = y0 + Inches(0.30) + ri * row_h
        bg = GREY_LT if ri % 2 == 0 else WHITE
        _rect(sl, Inches(0.25), y, Inches(12.85), row_h, fill=bg)
        _rect(sl, Inches(0.25), y, Inches(0.12), row_h, fill=PURPLE)

        # Col A: dim name + weight
        _txbox(sl, COL_A_X + Inches(0.2), y + Inches(0.08), Inches(2.1), Inches(0.26),
               dim, size=10, bold=True, color=DARK)
        _txbox(sl, COL_A_X + Inches(0.2), y + Inches(0.36), Inches(0.8), Inches(0.22),
               pct, size=11, bold=True, color=PURPLE)

        # Col B: what it measures
        _txbox(sl, COL_B_X + Inches(0.08), y + Inches(0.06),
               COL_B_W - Inches(0.14), row_h - Inches(0.1),
               what, size=7, color=DARK)

        # Col C: chip lines
        cy = y + Inches(0.08)
        for chip_color, label in chip_lines:
            cy = _chip_row(sl, COL_C_X + Inches(0.12), cy, chip_color, label, size=8)

    _footer(sl, date, mode)
    return sl


def _slide_appendix_divider(prs, date, mode):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, fill=PURPLE)
    _txbox(sl, Inches(1.0), Inches(2.8), Inches(11), Inches(1.2),
           "APPENDIX", size=48, bold=True, color=WHITE)
    _txbox(sl, Inches(1.0), Inches(4.1), Inches(11), Inches(0.5),
           "Detailed Scoring Methodology & Calculation Examples",
           size=18, color=RGBColor(0xE9, 0xD5, 0xFF))
    # Footer without mode/date branding
    _txbox(sl, Inches(0.4), H - Inches(0.4), Inches(12), Inches(0.3),
           "Accenture  ·SPM Readiness Assessment  ·  AS-IS Profile",
           size=9, color=RGBColor(0xE9, 0xD5, 0xFF))
    return sl


def _slide_appendix_methodology(prs, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Appendix A — Scoring Methodology",
                 "Complete reference: formula, weights, thresholds, and not-collected rules")

    # ── Section 1: Formula ────────────────────────────────────────────────────
    _txbox(sl, Inches(0.3), Inches(1.38), Inches(12.7), Inches(0.28),
           "SCORING FORMULA", size=9, bold=True, color=GREY_TEXT)
    _rect(sl, Inches(0.3), Inches(1.7), Inches(12.7), Inches(0.72), fill=GREY_LT, line=GREY_MID)
    _txbox(sl, Inches(0.5), Inches(1.78), Inches(12.3), Inches(0.56),
           "Module Score  =  Σ (dimension_score × dimension_weight)  ÷  Σ (weights of scored dimensions only)\n"
           "Overall Score  =  simple average of all module scores where module_score is not null (not_collected modules excluded)",
           size=10, color=DARK)

    # ── Section 2: Weights & thresholds side by side ─────────────────────────
    _txbox(sl, Inches(0.3), Inches(2.55), Inches(6.0), Inches(0.28),
           "DIMENSION WEIGHTS", size=9, bold=True, color=GREY_TEXT)
    _txbox(sl, Inches(6.9), Inches(2.55), Inches(6.0), Inches(0.28),
           "RAG THRESHOLDS", size=9, bold=True, color=GREY_TEXT)

    dim_rows = [
        ("Activation",        "20%"),
        ("Data Volume",       "20%"),
        ("Data Completeness", "25%"),
        ("Process Adoption",  "25%"),
        ("Integration",       "10%"),
        ("TOTAL",             "100%"),
    ]
    row_h = Inches(0.38)
    for i, (d, w) in enumerate(dim_rows):
        y = Inches(2.88) + i * row_h
        bg = PURPLE if d == "TOTAL" else (GREY_LT if i % 2 == 0 else WHITE)
        fc = WHITE if d == "TOTAL" else DARK
        _rect(sl, Inches(0.3), y, Inches(6.0), row_h, fill=bg)
        _txbox(sl, Inches(0.45), y + Inches(0.06), Inches(4.0), Inches(0.26), d, size=10, bold=(d=="TOTAL"), color=fc)
        _txbox(sl, Inches(5.5), y + Inches(0.06), Inches(0.7), Inches(0.26), w, size=10, bold=True, color=fc, align=PP_ALIGN.RIGHT)

    rag_rows = [
        ("Green",         "Score ≥ 70%",   RAG_RGB["green"]),
        ("Amber",         "40% ≤ Score < 70%", RAG_RGB["amber"]),
        ("Red",           "Score < 40%",   RAG_RGB["red"]),
        ("Not Collected", "Dimension excluded from average", RGBColor(0x9C, 0xA3, 0xAF)),
    ]
    for i, (label, meaning, col) in enumerate(rag_rows):
        y = Inches(2.88) + i * Inches(0.56)
        _rect(sl, Inches(6.9), y, Inches(1.2), Inches(0.46), fill=col)
        _txbox(sl, Inches(6.95), y + Inches(0.1), Inches(1.1), Inches(0.26),
               label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _rect(sl, Inches(8.15), y, Inches(4.8), Inches(0.46), fill=GREY_LT)
        _txbox(sl, Inches(8.25), y + Inches(0.1), Inches(4.6), Inches(0.26),
               meaning, size=10, color=DARK)

    # ── Section 3: Not-collected rules ────────────────────────────────────────
    _txbox(sl, Inches(0.3), Inches(5.22), Inches(12.7), Inches(0.28),
           "NOT-COLLECTED RULES", size=9, bold=True, color=GREY_TEXT)
    rules = [
        "A dimension scores not-collected (—) when its source data is absent or the plugin is inactive.",
        "Not-collected dimensions are excluded from the module score denominator — they do not count as zero.",
        "A module with all dimensions not-collected shows as not-collected overall and is excluded from the overall score.",
        "Financial completeness and integration are scored against PPM projects even when the financial plugin is off — "
        "because the gap is measurable (0 of N projects have cost/budget plans).",
    ]
    for i, rule in enumerate(rules):
        y = Inches(5.55) + i * Inches(0.28)
        _txbox(sl, Inches(0.5), y, Inches(12.5), Inches(0.26),
               f"•  {rule}", size=9, color=DARK)

    _footer(sl, date, mode)
    return sl


def _slide_appendix_example(prs, metrics, scores, date, mode):
    """Walk through one fully-scored and one partially-scored module as a worked example."""
    sl = _blank(prs)
    _header_band(sl, "Appendix B — Worked Calculation Examples",
                 "Step-by-step derivation of module scores from raw metrics")

    mods = metrics.get("modules", {})

    # Pick the two most instructive modules: best fully-scored + PPM
    fully_scored = [(k, scores[k]) for k in _MODULE_LABELS
                    if scores.get(k, {}).get("module_score") is not None]
    if not fully_scored:
        _txbox(sl, Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
               "No scored modules available.", size=12, color=GREY_TEXT)
        _footer(sl, date, mode)
        return sl

    # Use the first two scored modules (or just one if only one scored)
    examples = fully_scored[:2]
    col_x = [Inches(0.25), Inches(6.8)]

    for col_i, (mod_key, mod_scores) in enumerate(examples):
        x   = col_x[col_i]
        cw  = Inches(6.4)
        mod = mods.get(mod_key, {})
        ms  = mod_scores.get("module_score")
        rag = _rag_label(ms)
        col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))

        # Module header
        _rect(sl, x, Inches(1.38), cw, Inches(0.38), fill=PURPLE)
        _txbox(sl, x + Inches(0.12), Inches(1.44), Inches(4.0), Inches(0.26),
               _MODULE_LABELS[mod_key], size=11, bold=True, color=WHITE)
        _txbox(sl, x + Inches(4.2), Inches(1.44), Inches(2.0), Inches(0.26),
               f"Score = {ms}%  [{rag.upper()}]", size=10, bold=True,
               color=WHITE, align=PP_ALIGN.RIGHT)

        # Column headers
        hdrs = ["Dimension", "Weight", "Raw Input", "Score", "Weighted"]
        cws  = [Inches(1.5), Inches(0.55), Inches(2.3), Inches(0.65), Inches(0.8)]
        y = Inches(1.8)
        hdr_h = Inches(0.3)
        _rect(sl, x, y, cw, hdr_h, fill=GREY_MID)
        cx = x
        for hdr, w in zip(hdrs, cws):
            _txbox(sl, cx + Inches(0.04), y + Inches(0.05), w - Inches(0.06), Inches(0.2),
                   hdr, size=8, bold=True, color=DARK)
            cx += w

        # Dimension rows
        dims_info = [
            ("Activation",        "20%", 0.20, "activation"),
            ("Data Volume",       "20%", 0.20, "data_volume"),
            ("Data Completeness", "25%", 0.25, "data_completeness"),
            ("Process Adoption",  "25%", 0.25, "process_adoption"),
            ("Integration",       "10%", 0.10, "integration"),
        ]
        row_h = Inches(0.38)
        y += hdr_h
        total_w_used = 0.0
        total_weighted = 0.0

        for ri, (dim_label, wt_str, wt, dim_key) in enumerate(dims_info):
            score = mod_scores.get(dim_key)
            bg = GREY_LT if ri % 2 == 0 else WHITE
            _rect(sl, x, y, cw, row_h, fill=bg)

            # Raw input description
            raw_desc = _raw_input_desc(mod_key, dim_key, mod, score)
            weighted_str = "—"
            if score is not None:
                total_w_used += wt
                total_weighted += wt * score
                weighted_str = f"{round(wt * score, 1)}"

            score_str = f"{score}%" if score is not None else "—"
            row_vals = [dim_label, wt_str, raw_desc, score_str, weighted_str]
            cx = x
            for vi, (val, w) in enumerate(zip(row_vals, cws)):
                tc = RAG_RGB.get(_rag_label(score), RGBColor(0x9C, 0xA3, 0xAF)) if vi == 3 and score is not None else DARK
                _txbox(sl, cx + Inches(0.04), y + Inches(0.08), w - Inches(0.06), Inches(0.22),
                       val, size=8, color=tc, bold=(vi == 3 and score is not None))
                cx += w
            y += row_h

        # Totals row
        denom = round(total_w_used, 2)
        final  = round(total_weighted / total_w_used) if total_w_used else 0
        _rect(sl, x, y, cw, row_h, fill=PURPLE)
        _txbox(sl, x + Inches(0.12), y + Inches(0.08), Inches(4.0), Inches(0.24),
               f"Total weight used: {denom:.2f}  ÷  Weighted sum: {round(total_weighted, 1)}",
               size=8, color=WHITE)
        _txbox(sl, x + Inches(4.8), y + Inches(0.08), Inches(1.5), Inches(0.24),
               f"= {final}%", size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    _footer(sl, date, mode)
    return sl


def _raw_input_desc(mod_key, dim_key, mod, score):
    """One-line description of what raw metric(s) drove this dimension score."""
    inputs = _MODULE_INPUTS.get(mod_key, [])
    relevant = [lbl for lbl, mk, is_pct, inv, d_lbl in inputs
                if d_lbl.lower().replace(" ", "_") in dim_key or dim_key in d_lbl.lower()]
    if not relevant:
        relevant = [lbl for lbl, mk, _, _, d_lbl in inputs][:1]

    # Build concise value strings from known metrics
    if dim_key == "activation":
        return "Plugin active" if mod.get("plugin_active") else "Plugin inactive"
    if dim_key == "data_volume":
        total_key = "total_ci" if mod_key == "csdm" else ("total_stories" if mod_key == "agile" else "total")
        n = mod.get(total_key)
        return f"{n} records" if n is not None else "0 records"
    vals = []
    for lbl, metric_key, is_pct, invert, d_lbl in inputs:
        if d_lbl.lower().replace(" ", "_") not in dim_key and dim_key not in d_lbl.lower():
            continue
        raw = mod.get(metric_key)
        if raw is not None:
            v = round(100.0 - float(raw)) if invert else round(float(raw))
            vals.append(f"{lbl}: {v}%")
    if vals:
        return "  ·  ".join(vals[:2])
    return "—" if score is None else f"derived score: {score}%"


def _slide_appendix_tables(prs, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Appendix C — Source Tables",
                 "All ServiceNow tables queried during this assessment, grouped by module")

    _GREY_NC = RGBColor(0x9C, 0xA3, 0xAF)
    HDR_H, ROW_H, GAP = Inches(0.26), Inches(0.22), Inches(0.05)

    sections_left = [
        ("Demand Management", [
            ("pm_demand",           "Volume · completeness · approval rate · project linkage"),
            ("pm_demand_category",  "Demand categorisation breakdown"),
        ]),
        ("Project Portfolio (PPM)", [
            ("pm_project",          "Volume · completeness · status reports · approval rate"),
            ("pm_project_task",     "Task completion rates"),
            ("pm_program",          "Program linkage (integration)"),
        ]),
        ("Resource Management", [
            ("pm_resource_plan",       "Volume · named resource fill · project linkage"),
            ("pm_resource_allocation", "Actual vs planned hours (completeness)"),
        ]),
        ("Financial Management", [
            ("pm_project_financials", "Financial coverage per project (integration)"),
            ("pm_cost_plan",          "Cost plan linkage per project"),
            ("pm_budget_plan",        "Budget plan linkage per project"),
        ]),
        ("Agile Development", [
            ("rm_story",  "Volume · sprint assignment · team assignment"),
            ("rm_sprint", "Completed sprint count (process adoption)"),
            ("rm_team",   "Team reference for story / sprint assignment"),
        ]),
    ]

    sections_right = [
        ("Application Portfolio (APM)", [
            ("apm_appl_now",       "Volume · lifecycle stage · owner (completeness)"),
            ("apm_appl_lifecycle", "Lifecycle stage details"),
        ]),
        ("Innovation Management", [
            ("innovation_idea",      "Volume · owner fill rate · conversion rate"),
            ("innovation_challenge", "Challenge linkage"),
        ]),
        ("CSDM / CMDB Health", [
            ("cmdb_ci",          "Volume · field completeness · % discovered (adoption)"),
            ("cmdb_ci_service",  "Service count · service owner coverage"),
            ("cmdb_rel_ci",      "Relationship density — total rels ÷ total CIs (integration)"),
        ]),
        ("Governance Overlays", [
            ("timesheet_period",     "Active timesheet periods"),
            ("timesheet_entry",      "Entries logged per period (Resource adoption)"),
            ("pm_project_status",    "Status reports filed per project (PPM adoption)"),
            ("sysapproval_approver", "Approval records for projects and demands"),
        ]),
        ("Scoring Models", [
            ("pm_scoring_criterion", "Scoring criteria defined on instance"),
            ("pm_portfolio_score",   "Records scored — demand_scored signal (Demand adoption)"),
        ]),
        ("Performance Analytics", [
            ("pa_scorecard", "Active PA scorecards present on instance"),
        ]),
        ("Sidecars / Computed Inputs", [
            ("_sidecar_spm_adoption",     "Plugin active flags → feeds Activation dimension"),
            ("_sidecar_portfolio_health", "Stale record rates → portfolio health metrics"),
        ]),
    ]

    def _render_sections(slide, sections, col_x, col_w):
        y = Inches(1.38)
        tbl_w  = col_w - Inches(0.08)
        name_w = Inches(2.15)
        desc_w = tbl_w - name_w - Inches(0.06)

        for sec_label, tables in sections:
            # Section header
            _rect(slide, col_x, y, tbl_w, HDR_H, fill=PURPLE)
            _txbox(slide, col_x + Inches(0.08), y + Inches(0.05),
                   tbl_w - Inches(0.1), Inches(0.18),
                   sec_label, size=8, bold=True, color=WHITE)
            y += HDR_H

            for ri, (tbl, desc) in enumerate(tables):
                bg = GREY_LT if ri % 2 == 0 else WHITE
                _rect(slide, col_x, y, tbl_w, ROW_H, fill=bg)
                # table name — bold, left
                _txbox(slide, col_x + Inches(0.08), y + Inches(0.04),
                       name_w, Inches(0.16), tbl, size=8, bold=True, color=DARK)
                # description — grey, right
                _txbox(slide, col_x + name_w + Inches(0.1), y + Inches(0.04),
                       desc_w, Inches(0.16), desc, size=7, color=GREY_TEXT)
                y += ROW_H

            y += GAP

    _render_sections(sl, sections_left,  Inches(0.25), Inches(6.35))
    _render_sections(sl, sections_right, Inches(6.73), Inches(6.35))

    # Divider line between columns
    _rect(sl, Inches(6.67), Inches(1.38), Inches(0.04), Inches(5.9), fill=GREY_MID)

    # Bottom note
    total = sum(len(t) for _, t in sections_left) + sum(len(t) for _, t in sections_right)
    _txbox(sl, Inches(0.25), Inches(7.08), Inches(12.85), Inches(0.22),
           f"{total} tables total  ·  Sidecars are JSON files computed from instance data and "
           "stored alongside collector outputs — they are not direct ServiceNow table queries.",
           size=7, color=GREY_TEXT, italic=True)

    _footer(sl, date, mode)
    return sl


def _slide_summary(prs, metrics, scores, findings, overall, date, mode):
    sl = _blank(prs)
    _header_band(sl, "Assessment Summary",
                 "Key signals from the AS-IS readiness profile")

    # ── Left panel: overall score + RAG breakdown ─────────────────────────────
    _rect(sl, Inches(0.3), Inches(1.35), Inches(3.8), Inches(5.7),
          fill=GREY_LT, line=GREY_MID)

    rag_o = _rag_label(overall)
    col_o = RAG_RGB.get(rag_o, RGBColor(0x9C, 0xA3, 0xAF))
    _txbox(sl, Inches(0.5), Inches(1.5), Inches(3.4), Inches(0.35),
           "OVERALL SCORE", size=9, bold=True, color=GREY_TEXT)
    _txbox(sl, Inches(0.45), Inches(1.8), Inches(3.8), Inches(1.4),
           f"{overall}%" if overall is not None else "—",
           size=64, bold=True, color=col_o)

    # RAG breakdown counts
    rag_counts = {"green": 0, "amber": 0, "red": 0, "not_collected": 0}
    for key in _MODULE_LABELS:
        r = scores.get(key, {}).get("rag", "not_collected")
        rag_counts[r] = rag_counts.get(r, 0) + 1

    _txbox(sl, Inches(0.5), Inches(3.1), Inches(3.4), Inches(0.3),
           "MODULE BREAKDOWN", size=9, bold=True, color=GREY_TEXT)

    breakdown = [
        ("Green",         rag_counts["green"],         RAG_RGB["green"]),
        ("Amber",         rag_counts["amber"],          RAG_RGB["amber"]),
        ("Red",           rag_counts["red"],            RAG_RGB["red"]),
        ("Not Collected", rag_counts["not_collected"],  RGBColor(0x9C, 0xA3, 0xAF)),
    ]
    for i, (label, count, col) in enumerate(breakdown):
        y = Inches(3.5) + i * Inches(0.62)
        _rect(sl, Inches(0.5), y, Inches(0.28), Inches(0.28), fill=col)
        _txbox(sl, Inches(0.9), y - Inches(0.02), Inches(1.6), Inches(0.32),
               label, size=10, color=DARK)
        _txbox(sl, Inches(2.6), y - Inches(0.02), Inches(1.2), Inches(0.32),
               str(count), size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # Plugin activation summary
    mods = metrics.get("modules", {})
    active = sum(1 for k in _MODULE_LABELS if mods.get(k, {}).get("plugin_active"))
    _txbox(sl, Inches(0.5), Inches(6.1), Inches(3.4), Inches(0.28),
           f"{active} of {len(_MODULE_LABELS)} modules active",
           size=9, color=GREY_TEXT, italic=True)

    # ── Right panel: top findings ──────────────────────────────────────────────
    _txbox(sl, Inches(4.4), Inches(1.35), Inches(8.6), Inches(0.32),
           "KEY SIGNALS", size=9, bold=True, color=GREY_TEXT)

    top = findings[:5]
    row_h = Inches(0.92)
    for i, f in enumerate(top):
        y   = Inches(1.72) + i * (row_h + Inches(0.06))
        rag = f.get("rag", "not_collected")
        col = RAG_RGB.get(rag, RGBColor(0x9C, 0xA3, 0xAF))

        _rect(sl, Inches(4.4), y, Inches(8.6), row_h, fill=GREY_LT, line=GREY_MID)
        _rect(sl, Inches(4.4), y, Inches(0.22), row_h, fill=col)

        # Finding ID + module
        _txbox(sl, Inches(4.72), y + Inches(0.06), Inches(2.5), Inches(0.28),
               f"{f.get('id','')}  ·  {f.get('module_label', f.get('module',''))}",
               size=8, bold=True, color=col)
        # Observation
        _txbox(sl, Inches(4.72), y + Inches(0.32), Inches(8.0), Inches(0.36),
               f.get("observation", ""), size=9, color=DARK)
        # Significance
        sig = f.get("significance", "")
        if sig:
            _txbox(sl, Inches(4.72), y + Inches(0.62), Inches(8.0), Inches(0.26),
                   sig, size=8, color=GREY_TEXT, italic=True)

    _footer(sl, date, mode)
    return sl


# ── Public API ────────────────────────────────────────────────────────────────
def render_pptx(metrics, scores, findings, mode="rde"):
    ctx     = metrics.get("_context", {})
    client  = ctx.get("client", "Client").upper()
    date    = ctx.get("generated_on", "")
    mod_scores = {k: v.get("module_score") for k, v in scores.items()}
    overall    = _overall_score(mod_scores)

    prs = _prs()
    _slide_cover(prs, client, date, overall, mode)
    _slide_bar(prs, scores, overall, date, mode)
    _slide_scorecard(prs, scores, date, mode)
    _slide_scorecard_explainer(prs, metrics, scores, date, mode)
    _slide_scoring_basis(prs, metrics, scores, date, mode)
    _slide_scoring_basis_explainer(prs, metrics, scores, date, mode)
    _slide_governance(prs, metrics, date, mode)
    _slide_staleness(prs, metrics, date, mode)
    _slide_findings(prs, findings, date, mode)
    _slide_next_steps(prs, mode, findings, date)
    _slide_summary(prs, metrics, scores, findings, overall, date, mode)
    _slide_appendix_divider(prs, date, mode)
    _slide_appendix_methodology(prs, date, mode)
    _slide_appendix_example(prs, metrics, scores, date, mode)
    _slide_appendix_tables(prs, date, mode)
    return prs


def write_pptx(metrics, scores, findings, out_dir, mode="rde"):
    os.makedirs(out_dir, exist_ok=True)
    prs  = render_pptx(metrics, scores, findings, mode=mode)
    path = os.path.join(out_dir, "spm-leadership-deck.pptx")
    prs.save(path)
    return path


def main(argv=None):
    ap = argparse.ArgumentParser(description="Render SPM leadership PowerPoint deck")
    ap.add_argument("--metrics", required=True)
    ap.add_argument("--out",     required=True)
    ap.add_argument("--mode",    default="rde", choices=["rde", "fde"])
    args = ap.parse_args(argv)

    with open(args.metrics, encoding="utf-8") as f:
        metrics = json.load(f)

    from scripts.scoring  import score_all, enrich_coverage_matrix
    from scripts.metrics  import _coverage_matrix
    from scripts.findings import generate_findings

    if not metrics.get("coverage_matrix"):
        metrics["coverage_matrix"] = _coverage_matrix(metrics.get("modules", {}))
    scores   = score_all(metrics)
    metrics["coverage_matrix"] = enrich_coverage_matrix(metrics, scores)
    findings = generate_findings(metrics, scores)

    path = write_pptx(metrics, scores, findings, args.out, mode=args.mode)
    print(f"PowerPoint written: {path}")


if __name__ == "__main__":
    main()
